import os
import secrets
import shutil
import pandas as pd
from datetime import timedelta
from flask import Flask, render_template, request, send_file, jsonify, Response, stream_with_context, session, redirect, url_for
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io
import json
import requests
import tempfile
from PIL import Image
from collections import defaultdict
import concurrent.futures
import base64

import config
from auth import (
    auth_error_message,
    clear_user_session,
    current_user_email,
    current_user_id,
    get_supabase,
    is_authenticated,
    login_required,
    set_user_session,
    user_upload_dir,
)
from subscription import (
    apply_subscription_webhook,
    checkout_url_for_current_user,
    ensure_profile,
    pro_required,
    user_has_pro_access,
    verify_lemonsqueezy_signature,
)

app = Flask(__name__)
app.secret_key = config.SECRET_KEY
app.config["UPLOAD_ROOT"] = str(config.UPLOAD_ROOT)
app.config["ALLOWED_EXTENSIONS"] = config.ALLOWED_EXTENSIONS
app.config["PERMANENT_SESSION_LIFETIME"] = timedelta(days=7)
app.config["SESSION_COOKIE_HTTPONLY"] = True
app.config["SESSION_COOKIE_SAMESITE"] = "Lax"


@app.context_processor
def inject_auth():
    is_pro = False
    user_plan = "free"
    if is_authenticated():
        try:
            is_pro = user_has_pro_access()
            user_plan = "pro" if is_pro else "free"
        except Exception:
            is_pro = False
            user_plan = "free"
    return {
        "is_authenticated": is_authenticated(),
        "current_user_email": current_user_email(),
        "current_user_id": current_user_id(),
        "is_pro": is_pro,
        "user_plan": user_plan,
    }

# Resim indirme fonksiyonu (Tekli - Thread içinde kullanılacak)
def download_image_single(url, session):
    try:
        if not url or pd.isna(url):
            return None, None
        
        # URL string değilse atla
        if not isinstance(url, str):
            return None, None

        response = session.get(url, timeout=10)
        if response.status_code == 200:
            return url, response.content
        else:
            return url, None
    except:
        return url, None

# Toplu ve Paralel Resim İndirme Fonksiyonu
def download_images_parallel(urls, max_workers=20):
    """
    Verilen URL listesindeki resimleri paralel olarak indirir.
    Geriye {url: resim_bytes} şeklinde bir sözlük döndürür.
    """
    image_cache = {}
    # Benzersiz URL'leri al
    unique_urls = list(set([url for url in urls if pd.notna(url) and isinstance(url, str) and str(url).strip() != '']))
    
    # Session oluştur
    session = requests.Session()
    
    print(f"Toplam {len(unique_urls)} adet resim paralel indiriliyor...")
    
    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_to_url = {executor.submit(download_image_single, url, session): url for url in unique_urls}
        
        for future in concurrent.futures.as_completed(future_to_url):
            url, content = future.result()
            if content:
                image_cache[url] = content
                
    return image_cache

# GÜNCELLENDİ: Slayt arka planını "Resim Şekli" olarak ekleme fonksiyonu
def set_slide_background(slide, image_path, prs):
    """
    Verilen resmi slaytın arka planı olarak ayarlar.
    Garanti görünmesi için resmi tam ekran boyutunda ve en arkaya (z-order 0) ekler.
    """
    if not image_path:
        return
    
    try:
        # Slayt boyutlarını al
        left = top = Inches(0)
        width = prs.slide_width
        height = prs.slide_height
        
        # Resmi ekle (Slayt yeni oluşturulduğunda ilk eklenen nesne en altta olur)
        slide.shapes.add_picture(image_path, left, top, width, height)
        
    except Exception as e:
        print(f"Arka plan ayarlanamadı: {e}")

def add_image_to_slide(slide, img_url, presentation_type, image_index, config, image_cache=None):
    """Resmi slayta ekle"""
    if presentation_type == 'single':
        image_config = config.get('image', {})
        left = Inches(image_config.get('left', 1))
        top = Inches(image_config.get('top', 1.5))
        width = Inches(image_config.get('width', 8))
        height = Inches(image_config.get('height', 6))
    elif presentation_type == 'double':
        if image_index == 0:
            image_config = config.get('image1', {})
        else:
            image_config = config.get('image2', {})
        left = Inches(image_config.get('left', 0.5))
        top = Inches(image_config.get('top', 1.5))
        width = Inches(image_config.get('width', 4))
        height = Inches(image_config.get('height', 3))
    else:  # triple
        if image_index == 0:
            image_config = config.get('image1', {})
        elif image_index == 1:
            image_config = config.get('image2', {})
        else:
            image_config = config.get('image3', {})
        left = Inches(image_config.get('left', 0.3))
        top = Inches(image_config.get('top', 1.5))
        width = Inches(image_config.get('width', 2.5))
        height = Inches(image_config.get('height', 2))
    
    try:
        image_content = None
        if image_cache and img_url in image_cache:
            image_content = image_cache[img_url]
        elif not image_cache:
            try:
                response = requests.get(img_url, timeout=10)
                if response.status_code == 200:
                    image_content = response.content
            except:
                image_content = None

        if image_content:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as tmp_file:
                tmp_file.write(image_content)
                tmp_path = tmp_file.name
            
            with Image.open(tmp_path) as img:
                img.thumbnail((int(width.inches * 150), int(height.inches * 150)), Image.Resampling.LANCZOS)
                img.save(tmp_path, 'JPEG', quality=95)
            
            slide.shapes.add_picture(tmp_path, left, top, width, height)
            os.unlink(tmp_path)
        else:
            raise Exception("Empty image content")

    except Exception as e:
        placeholder = slide.shapes.add_shape(1, left, top, width, height)
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(230, 230, 230)
        textbox = slide.shapes.add_textbox(left, top + height, width, Inches(0.5))
        textframe = textbox.text_frame
        textframe.text = "Image failed to load"
        textframe.paragraphs[0].font.size = Pt(10)
        textframe.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        textframe.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_title_to_slide(slide, title_text, presentation_type, title_index, config, bold=True, word_wrap=True):
    """Başlığı slayta ekle"""
    if presentation_type == 'single':
        title_config = config.get('title', {})
        title_left = Inches(title_config.get('left', 1))
        title_top = Inches(title_config.get('top', 1))
        title_width = Inches(title_config.get('width', 8))
        title_height = Inches(title_config.get('height', 0.4))
        title_font_size = title_config.get('fontSize', 14)
    elif presentation_type == 'double':
        if title_index == 0:
            title_config = config.get('title1', {})
        else:
            title_config = config.get('title2', {})
        title_left = Inches(title_config.get('left', 0.5))
        title_top = Inches(title_config.get('top', 1))
        title_width = Inches(title_config.get('width', 4))
        title_height = Inches(title_config.get('height', 0.4))
        title_font_size = title_config.get('fontSize', 14)
    else:  # triple
        if title_index == 0:
            title_config = config.get('title1', {})
        elif title_index == 1:
            title_config = config.get('title2', {})
        else:
            title_config = config.get('title3', {})
        title_left = Inches(title_config.get('left', 0.3))
        title_top = Inches(title_config.get('top', 1))
        title_width = Inches(title_config.get('width', 2.5))
        title_height = Inches(title_config.get('height', 0.4))
        title_font_size = title_config.get('fontSize', 12)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = str(title_text) if pd.notna(title_text) else ""
    title_frame.word_wrap = word_wrap
    title_frame.paragraphs[0].font.bold = bold
    title_frame.paragraphs[0].font.size = Pt(title_font_size)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def _get_cell_value_as_string(row_data, col_name):
    """Satır verisinden güvenli string değer üretir (NaN/None -> '')."""
    try:
        if hasattr(row_data, 'get') and col_name in row_data:
            val = row_data[col_name]
        else:
            val = getattr(row_data, col_name, None)
    except Exception:
        val = None

    if val is None or (isinstance(val, float) and pd.isna(val)) or pd.isna(val):
        return ""
    return str(val)

def build_caption_text(row_data, caption_columns):
    """
    Seçilen sütunlara göre başlık metnini üretir.
    Format:
      Sütun1
      Değer1
      Sütun2
      Değer2
    """
    # row_data genelde pandas Series gelir; Series'i boolean context'te kullanmak "ambiguous" hatası üretir.
    if row_data is None or not caption_columns:
        return ""

    values = []
    for col in caption_columns:
        value_str = _get_cell_value_as_string(row_data, col).strip()
        if value_str != "":
            values.append(value_str)

    # İstenen format: sadece değer(ler). Birden fazla sütun seçilirse alt alta.
    return "\n".join(values)

def add_caption_to_slide(slide, row_data, caption_columns, presentation_type, title_index, config, bold=True, word_wrap=True):
    """
    Satır verisinden çok satırlı fotoğraf başlığı ekler.
    Başlık kutusunun konum/ölçüsü `title/title1/title2/title3` config'lerinden okunur.
    """
    caption_text = build_caption_text(row_data, caption_columns)
    if not caption_text:
        return

    # Konum/ölçü: mevcut title config'lerini kullan
    if presentation_type == 'single':
        title_config = config.get('title', {})
        title_left = Inches(title_config.get('left', 1))
        title_top = Inches(title_config.get('top', 1))
        title_width = Inches(title_config.get('width', 8))
        title_height = Inches(title_config.get('height', 0.4))
        title_font_size = title_config.get('fontSize', 14)
    elif presentation_type == 'double':
        title_config = config.get('title1', {}) if title_index == 0 else config.get('title2', {})
        title_left = Inches(title_config.get('left', 0.5))
        title_top = Inches(title_config.get('top', 1))
        title_width = Inches(title_config.get('width', 4))
        title_height = Inches(title_config.get('height', 0.4))
        title_font_size = title_config.get('fontSize', 14)
    else:  # triple
        if title_index == 0:
            title_config = config.get('title1', {})
        elif title_index == 1:
            title_config = config.get('title2', {})
        else:
            title_config = config.get('title3', {})
        title_left = Inches(title_config.get('left', 0.3))
        title_top = Inches(title_config.get('top', 1))
        title_width = Inches(title_config.get('width', 2.5))
        title_height = Inches(title_config.get('height', 0.4))
        title_font_size = title_config.get('fontSize', 12)

    text_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = text_box.text_frame
    tf.clear()
    tf.word_wrap = word_wrap
    tf.vertical_anchor = 1  # TOP

    # caption_text: sadece değerler (çoklu seçimde alt alta)
    lines = [ln for ln in caption_text.split("\n") if ln.strip() != ""]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = ""
        run = p.add_run()
        run.text = line
        run.font.size = Pt(title_font_size)
        run.font.name = "Arial"
        run.font.bold = bold
        run.font.color.rgb = RGBColor(33, 37, 41)
        p.alignment = PP_ALIGN.CENTER

def add_labels_for_image(slide, row_data, label_columns, label_config_key, add_label_border, config):
    """Etiketleri slayta ekle"""
    label_config = config.get(label_config_key, {})
    left = Inches(label_config.get('left', 1))
    top = Inches(label_config.get('top', 6.5))
    width = Inches(label_config.get('width', 8))
    height = Inches(label_config.get('height', 0.8))
    font_size = label_config.get('fontSize', 12)
    
    label_text = ""
    for col in label_columns:
        if hasattr(row_data, 'get') and col in row_data and pd.notna(row_data[col]):
            label_text += f"{col}: {row_data[col]}\n"
        elif hasattr(row_data, col) and pd.notna(getattr(row_data, col, None)):
            label_text += f"{col}: {getattr(row_data, col)}\n"
    
    if label_text:
        if add_label_border:
            label_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            label_box.fill.solid()
            label_box.fill.fore_color.rgb = RGBColor(248, 249, 250)
            label_box.line.color.rgb = RGBColor(206, 212, 218)
            label_box.line.width = Pt(1)
            text_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.05), width - Inches(0.2), height - Inches(0.1))
        else:
            text_box = slide.shapes.add_textbox(left, top, width, height)
        
        text_frame = text_box.text_frame
        text_frame.text = label_text.strip()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = 1
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.name = "Arial"
            paragraph.font.color.rgb = RGBColor(33, 37, 41)
            text_parts = paragraph.text.split(":")
            if len(text_parts) > 1:
                paragraph.text = ""
                run = paragraph.add_run()
                run.text = text_parts[0] + ":"
                run.font.bold = True
                run.font.color.rgb = RGBColor(13, 110, 253)
                run = paragraph.add_run()
                run.text = text_parts[1]
                run.font.bold = False
                run.font.color.rgb = RGBColor(33, 37, 41)

@app.route('/register', methods=['GET', 'POST'])
def register():
    if is_authenticated():
        return redirect(url_for('dashboard'))
    error = None
    email = ''
    if request.method == 'POST':
        email = (request.form.get('email') or '').strip()
        password = request.form.get('password') or ''
        password2 = request.form.get('password2') or ''
        if not email or not password:
            error = 'Email and password are required.'
        elif password != password2:
            error = 'Passwords do not match.'
        elif len(password) < 6:
            error = 'Password must be at least 6 characters.'
        else:
            try:
                sb = get_supabase()
                result = sb.auth.sign_up({"email": email, "password": password})
                user = getattr(result, "user", None)
                sess = getattr(result, "session", None)
                if user and sess:
                    set_user_session(
                        user,
                        access_token=getattr(sess, "access_token", None),
                        refresh_token=getattr(sess, "refresh_token", None),
                    )
                    user_upload_dir()
                    try:
                        ensure_profile(str(user.id), getattr(user, "email", None) or email)
                    except Exception:
                        pass
                    return redirect(url_for('dashboard'))
                if user and not sess:
                    try:
                        ensure_profile(str(user.id), getattr(user, "email", None) or email)
                    except Exception:
                        pass
                    return redirect(url_for('login', message='Registration received. If email confirmation is required, please verify and then log in.'))
                error = 'Registration could not be completed. Please try again.'
            except Exception as e:
                error = auth_error_message(e)
    return render_template('register.html', error=error, email=email)


@app.route('/login', methods=['GET', 'POST'])
def login():
    if is_authenticated():
        return redirect(url_for('dashboard'))
    error = None
    email = ''
    next_url = request.args.get('next') or request.form.get('next') or ''
    message = request.args.get('message') or ''
    if request.method == 'POST':
        email = (request.form.get('email') or '').strip()
        password = request.form.get('password') or ''
        if not email or not password:
            error = 'Email and password are required.'
        else:
            try:
                sb = get_supabase()
                result = sb.auth.sign_in_with_password({"email": email, "password": password})
                user = getattr(result, "user", None)
                sess = getattr(result, "session", None)
                if not user or not sess:
                    error = 'Login failed. Please check your credentials.'
                else:
                    set_user_session(
                        user,
                        access_token=getattr(sess, "access_token", None),
                        refresh_token=getattr(sess, "refresh_token", None),
                    )
                    user_upload_dir()
                    try:
                        ensure_profile(str(user.id), getattr(user, "email", None) or email)
                    except Exception:
                        pass
                    if next_url and next_url.startswith('/') and not next_url.startswith('//'):
                        return redirect(next_url)
                    return redirect(url_for('dashboard'))
            except Exception as e:
                error = auth_error_message(e)
    return render_template('login.html', error=error, email=email, next=next_url, message=message)


@app.route('/logout', methods=['GET', 'POST'])
def logout():
    try:
        if session.get('access_token'):
            get_supabase().auth.sign_out()
    except Exception:
        pass
    clear_user_session()
    return redirect(url_for('login'))


@app.route('/pricing')
def pricing():
    reason = request.args.get('reason') or ''
    checkout_url = ''
    is_pro = False
    if is_authenticated():
        checkout_url = checkout_url_for_current_user()
        try:
            is_pro = user_has_pro_access()
        except Exception:
            is_pro = False
    return render_template(
        'pricing.html',
        reason=reason,
        checkout_url=checkout_url,
        is_pro=is_pro,
    )


@app.route('/api/webhook/lemonsqueezy', methods=['POST'])
def lemonsqueezy_webhook():
    """Lemon Squeezy abonelik webhook'u — auth gerektirmez; imza doğrulanır."""
    raw = request.get_data()
    signature = request.headers.get('X-Signature') or request.headers.get('x-signature')
    if not verify_lemonsqueezy_signature(raw, signature):
        return jsonify({'error': 'Invalid signature'}), 401
    try:
        payload = request.get_json(force=True, silent=False)
    except Exception:
        try:
            payload = json.loads(raw.decode('utf-8'))
        except Exception:
            return jsonify({'error': 'Invalid JSON'}), 400
    if not isinstance(payload, dict):
        return jsonify({'error': 'Invalid payload'}), 400
    try:
        result = apply_subscription_webhook(payload)
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    if not result.get('ok') and result.get('error') == 'profile_not_found':
        # Lemon Squeezy tekrar denemesin diye 200 + uyarı (veya 404)
        return jsonify(result), 200
    return jsonify(result), 200


@app.route('/')
def landing():
    """Public marketing homepage. Authenticated users go to the app dashboard."""
    if is_authenticated():
        return redirect(url_for('dashboard'))
    from datetime import datetime
    return render_template('index.html', current_year=datetime.now().year)


@app.route('/dashboard')
@login_required
def dashboard():
    """Signed-in home: choose Slide Generator or Photo Extractor."""
    return render_template('home.html')


@app.route('/sunum')
@login_required
def sunum_app():
    return render_template('sunum.html')


@app.route('/rapor-fotograf')
@login_required
def rapor_fotograf_app():
    return render_template('rapor_fotograf.html')


@app.route('/api/rapor-fotograf/count', methods=['POST'])
@login_required
def rapor_fotograf_count():
    data = request.json or {}
    filename = data.get('filename')
    options = data.get('options') or {}
    if not filename:
        return jsonify({'error': 'Filename is required'}), 400
    upload_dir = user_upload_dir()
    safe, filepath = resolve_user_upload(upload_dir, filename)
    if not filepath:
        return jsonify({'error': 'File not found. Please upload an Excel file first.'}), 400
    try:
        from rapor_fotograf_logic import count_photo_urls_for_options

        total = count_photo_urls_for_options(upload_dir, safe, options)
        return jsonify({'total': total})
    except Exception as e:
        print(f"[rapor-fotograf/count] {type(e).__name__}: {e}")
        return jsonify({'error': str(e)}), 400


@app.route('/api/rapor-fotograf/process', methods=['POST'])
@pro_required
def rapor_fotograf_process():
    data = request.json or {}
    filename = data.get('filename')
    options = data.get('options') or {}
    if not filename:
        return jsonify({'error': 'Filename is required'}), 400
    upload_dir = user_upload_dir()
    safe, filepath = resolve_user_upload(upload_dir, filename)
    if not filepath:
        return jsonify({'error': 'File not found. Please upload an Excel file first.'}), 400
    try:
        from rapor_fotograf_logic import run_photo_extraction
        result = run_photo_extraction(upload_dir, safe, options)
    except Exception as e:
        print(f"[rapor-fotograf/process] {type(e).__name__}: {e}")
        return jsonify({'error': str(e)}), 400

    if isinstance(result, dict) and result.get('mode') == 'zip':
        zp = result.get('zip_path')
        if zp and os.path.isfile(zp):
            resp = send_file(
                zp,
                as_attachment=True,
                download_name=os.path.basename(zp),
                mimetype='application/zip',
            )
            saved_n = result.get('saved', 0)
            resp.headers['X-Rapor-Image-Count'] = str(saved_n)
            return resp

    if isinstance(result, str) and os.path.isfile(result):
        return send_file(
            result,
            as_attachment=True,
            download_name=os.path.basename(result),
        )
    if isinstance(result, tuple) and len(result) >= 2:
        payload, download_name = result[0], result[1]
        if isinstance(payload, bytes):
            buf = io.BytesIO(payload)
            buf.seek(0)
            return send_file(
                buf,
                as_attachment=True,
                download_name=str(download_name),
                mimetype='application/octet-stream',
            )
    return jsonify({'error': 'Invalid response from photo extraction (expected file path or bytes payload)'}), 500


@app.route('/api/rapor-fotograf/process-stream', methods=['POST'])
@pro_required
def rapor_fotograf_process_stream():
    data = request.json or {}
    filename = data.get('filename')
    options = data.get('options') or {}
    if not filename:
        return jsonify({'error': 'Filename is required'}), 400
    upload_dir = user_upload_dir()
    safe, filepath = resolve_user_upload(upload_dir, filename)
    if not filepath:
        return jsonify({'error': 'File not found. Please upload an Excel file first.'}), 400

    def generate():
        from rapor_fotograf_logic import iter_photo_extraction

        try:
            for ev in iter_photo_extraction(upload_dir, safe, options):
                if ev.get('type') == 'done':
                    token = secrets.token_urlsafe(24)
                    zp = ev['zip_path']
                    pending_dir = os.path.join(upload_dir, '_rf_pending')
                    os.makedirs(pending_dir, exist_ok=True)
                    dest = os.path.join(pending_dir, token + '.zip')
                    shutil.move(zp, dest)
                    out = {k: v for k, v in ev.items() if k != 'zip_path'}
                    out['token'] = token
                    yield f"data: {json.dumps(out, ensure_ascii=False, default=str)}\n\n"
                else:
                    yield f"data: {json.dumps(ev, ensure_ascii=False, default=str)}\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'type': 'error', 'message': str(e)}, ensure_ascii=False, default=str)}\n\n"

    return Response(
        stream_with_context(generate()),
        mimetype='text/event-stream',
        headers={
            'Cache-Control': 'no-cache',
            'Connection': 'keep-alive',
            'X-Accel-Buffering': 'no',
        },
    )


@app.route('/api/rapor-fotograf/download/<token>')
@pro_required
def rapor_fotograf_download_token(token: str):
    if not token or any(c in token for c in ('/', '\\', '..')):
        return jsonify({'error': 'Invalid download link.'}), 400
    upload_dir = user_upload_dir()
    path = os.path.join(upload_dir, '_rf_pending', token + '.zip')
    if not os.path.isfile(path):
        return jsonify({'error': 'Invalid or expired download.'}), 404
    resp = send_file(
        path,
        as_attachment=True,
        download_name=os.path.basename(path),
        mimetype='application/zip',
    )

    @resp.call_on_close
    def _remove_zip() -> None:
        try:
            os.remove(path)
        except OSError:
            pass

    return resp


def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']


# Turkish / Latin extended → ASCII for safe filesystem names
_TR_TRANSLATE = str.maketrans({
    "ç": "c", "Ç": "C",
    "ğ": "g", "Ğ": "G",
    "ı": "i", "İ": "I",
    "ö": "o", "Ö": "O",
    "ş": "s", "Ş": "S",
    "ü": "u", "Ü": "U",
    "â": "a", "Â": "A",
    "î": "i", "Î": "I",
    "û": "u", "Û": "U",
})


def sanitize_upload_filename(original_name: str) -> str:
    """
    Normalize Turkish/special characters, then apply secure_filename.
    Falls back to a unique name if the result would be empty.
    """
    import unicodedata

    raw = (original_name or "").strip()
    if not raw:
        return f"upload_{secrets.token_hex(8)}.xlsx"

    base, ext = os.path.splitext(raw)
    ext = ext.lower() if ext else ""

    # Map Turkish letters first (before NFKD), then strip remaining accents
    base_mapped = base.translate(_TR_TRANSLATE)
    base_norm = unicodedata.normalize("NFKD", base_mapped)
    base_norm = "".join(c for c in base_norm if not unicodedata.combining(c))
    base_ascii = "".join(ch if ord(ch) < 128 else "_" for ch in base_norm)

    candidate = secure_filename(base_ascii + ext)
    if not candidate or candidate in (".", "..") or candidate.startswith("."):
        suffix = ext if ext in (".xlsx", ".xls") else ".xlsx"
        candidate = f"upload_{secrets.token_hex(8)}{suffix}"
    stem, sex = os.path.splitext(candidate)
    if not stem:
        candidate = f"upload_{secrets.token_hex(8)}{sex or '.xlsx'}"
    return candidate


def resolve_user_upload(upload_dir: str, filename: str) -> tuple[str, str] | tuple[None, None]:
    """Find an uploaded file under the user folder; prevent path traversal."""
    if not filename:
        return None, None
    name = os.path.basename(str(filename).strip())
    if not name or name in (".", "..") or "/" in name or "\\" in name:
        return None, None
    # Prefer exact name returned from upload, then sanitized variants
    candidates = []
    for c in (name, secure_filename(name), sanitize_upload_filename(name)):
        if c and c not in candidates:
            candidates.append(c)
    for c in candidates:
        path = os.path.join(upload_dir, c)
        if os.path.isfile(path):
            return c, path
    return None, None


def unique_upload_path(upload_dir: str, filename: str) -> tuple[str, str]:
    """Return (final_filename, full_path), appending _1, _2… if needed."""
    name = filename
    path = os.path.join(upload_dir, name)
    if not os.path.exists(path):
        return name, path
    stem, ext = os.path.splitext(name)
    n = 1
    while True:
        name = f"{stem}_{n}{ext}"
        path = os.path.join(upload_dir, name)
        if not os.path.exists(path):
            return name, path
        n += 1


@app.route('/upload', methods=['POST'])
@login_required
def upload_file():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file selected. Please choose an Excel file (.xlsx or .xls).'}), 400
        file = request.files['file']
        original_name = file.filename or ''
        if original_name == '':
            return jsonify({'error': 'No file selected. Please choose an Excel file (.xlsx or .xls).'}), 400

        if not allowed_file(original_name):
            return jsonify({
                'error': f'Invalid file format for "{original_name}". Only .xlsx and .xls are supported.'
            }), 400

        filename = sanitize_upload_filename(original_name)
        upload_dir = user_upload_dir()
        filename, filepath = unique_upload_path(upload_dir, filename)

        print(
            f"[upload] user={current_user_id()} original={original_name!r} "
            f"saved_as={filename!r} path={filepath!r}"
        )

        try:
            file.save(filepath)
        except Exception as e:
            print(f"[upload] save failed: {type(e).__name__}: {e}")
            return jsonify({
                'error': f'Could not save the uploaded file. Please try again. ({type(e).__name__})'
            }), 500

        if not os.path.isfile(filepath) or os.path.getsize(filepath) == 0:
            print(f"[upload] empty or missing after save: {filepath!r}")
            return jsonify({'error': 'Upload failed: the saved file is empty or missing.'}), 500

        try:
            excel_file = pd.ExcelFile(filepath)
            sheet_names = excel_file.sheet_names
            if not sheet_names:
                return jsonify({'error': 'The Excel file has no sheets.'}), 400
            sheets_data = {}
            for sheet in sheet_names:
                df = pd.read_excel(filepath, sheet_name=sheet)
                sheets_data[sheet] = [str(c) for c in df.columns.tolist()]
            print(f"[upload] ok sheets={sheet_names!r} cols={[len(v) for v in sheets_data.values()]}")
            return jsonify({
                'filename': filename,
                'original_filename': original_name,
                'sheets': sheet_names,
                'sheets_data': sheets_data,
            })
        except Exception as e:
            print(f"[upload] Excel read failed for {filepath!r}: {type(e).__name__}: {e}")
            try:
                os.remove(filepath)
            except OSError:
                pass
            return jsonify({
                'error': (
                    f'Could not read Excel file "{original_name}". '
                    f'Make sure it is a valid .xlsx/.xls workbook. Details: {type(e).__name__}: {e}'
                )
            }), 500
    except Exception as e:
        print(f"[upload] unexpected error: {type(e).__name__}: {e}")
        return jsonify({
            'error': f'Upload failed unexpectedly ({type(e).__name__}: {e}). Please try again or rename the file.'
        }), 500


@app.route('/create_presentation', methods=['POST'])
@pro_required
def create_presentation():
    data = request.json
    filename = data.get('filename')
    selected_sheets = data.get('selected_sheets', [])
    presentation_type = data.get('presentation_type', 'single')
    selected_columns = data.get('selected_columns', [])
    label_columns = data.get('label_columns', [])
    use_column_titles = data.get('use_column_titles', False)
    caption_columns = data.get('caption_columns', [])
    add_label_border = data.get('add_label_border', False)
    label_count = data.get('label_count', 1)
    skip_single_images = data.get('skip_single_images', False)
    group_by_name = data.get('group_by_name', False)
    grouping_column = data.get('grouping_column', '')
    position_config = data.get('position_config', {})
    column_pairs = data.get('column_pairs', [])
    skip_empty_rows = data.get('skip_empty_rows', False)
    font_title_bold = data.get('font_title_bold', True)
    font_title_word_wrap = data.get('font_title_word_wrap', True)
    background_image_data = data.get('background_image_data')

    if not filename:
        return jsonify({'error': 'Filename is required'}), 400
    upload_dir = user_upload_dir()
    safe_name, filepath = resolve_user_upload(upload_dir, filename)
    if not filepath:
        return jsonify({'error': 'File not found. Please upload an Excel file first.'}), 400
    
    temp_bg_path = None
    try:
        if background_image_data and ',' in background_image_data:
            header, encoded = background_image_data.split(",", 1)
            bg_binary_data = base64.b64decode(encoded)
            with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as f:
                f.write(bg_binary_data)
                temp_bg_path = f.name

        prs = Presentation()
        # 16:9 Geniş Ekran
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide_layout = prs.slide_layouts[6]
        
        for sheet in selected_sheets:
            df = pd.read_excel(filepath, sheet_name=sheet)
            image_columns = [col for col in selected_columns if col in df.columns]
            all_urls_in_sheet = []
            for col in image_columns:
                all_urls_in_sheet.extend(df[col].tolist())
            if presentation_type == 'double' and column_pairs:
                for pair in column_pairs:
                    if pair.get('col1') and pair.get('col1') in df.columns: all_urls_in_sheet.extend(df[pair.get('col1')].tolist())
                    if pair.get('col2') and pair.get('col2') in df.columns: all_urls_in_sheet.extend(df[pair.get('col2')].tolist())
            
            image_cache = download_images_parallel(all_urls_in_sheet)
            
            if group_by_name and grouping_column and grouping_column in df.columns:
                grouped_data = defaultdict(list)
                for index, row in df.iterrows():
                    group_name = row[grouping_column] if pd.notna(row[grouping_column]) else "Other"
                    grouped_data[group_name].append(row)
                
                for group_name, rows in grouped_data.items():
                    all_image_urls = []
                    all_image_columns = []
                    all_image_row_data = []
                    for row in rows:
                        row_image_urls = []
                        row_image_columns = []
                        for col in image_columns:
                            if pd.notna(row[col]):
                                row_image_urls.append(row[col])
                                row_image_columns.append(col)
                        if skip_single_images and presentation_type != 'single':
                            required_images = 2 if presentation_type == 'double' else 3
                            if len(row_image_urls) < required_images: continue
                        all_image_urls.extend(row_image_urls)
                        all_image_columns.extend(row_image_columns)
                        all_image_row_data.extend([row] * len(row_image_urls))
                    
                    if not all_image_urls: continue
                    if presentation_type == 'single': images_per_slide = 1
                    elif presentation_type == 'double': images_per_slide = 2
                    else: images_per_slide = 3
                    
                    if presentation_type == 'triple' and skip_empty_rows:
                        filtered_image_urls = []
                        filtered_image_columns = []
                        filtered_image_row_data = []
                        for url, col, row_data in zip(all_image_urls, all_image_columns, all_image_row_data):
                            if pd.notna(url) and url != '':
                                filtered_image_urls.append(url)
                                filtered_image_columns.append(col)
                                filtered_image_row_data.append(row_data)
                        
                        for i in range(0, len(filtered_image_urls), images_per_slide):
                            slide = prs.slides.add_slide(slide_layout)
                            # GÜNCELLENDİ: Arka plan set edilirken prs parametresi eklendi
                            set_slide_background(slide, temp_bg_path, prs)
                            
                            current_images = filtered_image_urls[i:i+images_per_slide]
                            current_columns = filtered_image_columns[i:i+images_per_slide]
                            current_row_data_list = filtered_image_row_data[i:i+images_per_slide]
                            config = position_config.get(presentation_type, {})
                            for j, img_url in enumerate(current_images):
                                add_image_to_slide(slide, img_url, presentation_type, j, config, image_cache)
                                if caption_columns and j < len(current_row_data_list):
                                    add_caption_to_slide(slide, current_row_data_list[j], caption_columns, presentation_type, j, config, font_title_bold, font_title_word_wrap)
                                elif use_column_titles and j < len(current_columns):
                                    column_title = current_columns[j]
                                    add_title_to_slide(slide, column_title, presentation_type, j, config, font_title_bold, font_title_word_wrap)
                                if label_columns and j < len(current_row_data_list):
                                    if presentation_type == 'single':
                                        for k in range(label_count):
                                            current_row_data = current_row_data_list[j]
                                            label_config_key = f'label{k+1}'
                                            add_labels_for_image(slide, current_row_data, label_columns, label_config_key, add_label_border, config)
                                    elif (j + 1) <= label_count:
                                        current_row_data = current_row_data_list[j]
                                        label_config_key = f'label{j+1}'
                                        add_labels_for_image(slide, current_row_data, label_columns, label_config_key, add_label_border, config)
                    else:
                        for i in range(0, len(all_image_urls), images_per_slide):
                            slide = prs.slides.add_slide(slide_layout)
                            # GÜNCELLENDİ: Arka plan set edilirken prs parametresi eklendi
                            set_slide_background(slide, temp_bg_path, prs)
                            
                            current_images = all_image_urls[i:i+images_per_slide]
                            current_columns = all_image_columns[i:i+images_per_slide]
                            current_row_data_list = all_image_row_data[i:i+images_per_slide]
                            config = position_config.get(presentation_type, {})
                            for j, img_url in enumerate(current_images):
                                add_image_to_slide(slide, img_url, presentation_type, j, config, image_cache)
                                if caption_columns and j < len(current_row_data_list):
                                    add_caption_to_slide(slide, current_row_data_list[j], caption_columns, presentation_type, j, config, font_title_bold, font_title_word_wrap)
                                elif use_column_titles and j < len(current_columns):
                                    column_title = current_columns[j]
                                    add_title_to_slide(slide, column_title, presentation_type, j, config, font_title_bold, font_title_word_wrap)
                                if label_columns and j < len(current_row_data_list):
                                    if presentation_type == 'single':
                                        for k in range(label_count):
                                            current_row_data = current_row_data_list[j]
                                            label_config_key = f'label{k+1}'
                                            add_labels_for_image(slide, current_row_data, label_columns, label_config_key, add_label_border, config)
                                    elif (j + 1) <= label_count:
                                        current_row_data = current_row_data_list[j]
                                        label_config_key = f'label{j+1}'
                                        add_labels_for_image(slide, current_row_data, label_columns, label_config_key, add_label_border, config)
            else:
                if presentation_type == 'double' and column_pairs:
                    for index, row in df.iterrows():
                        for pair in column_pairs:
                            col1 = pair.get('col1'); col2 = pair.get('col2')
                            if not col1 or not col2: continue
                            img_url1 = row[col1] if pd.notna(row[col1]) else None
                            img_url2 = row[col2] if pd.notna(row[col2]) else None
                            if skip_single_images and (not img_url1 or not img_url2): continue
                            if not img_url1 and not img_url2: continue
                            
                            slide = prs.slides.add_slide(slide_layout)
                            # GÜNCELLENDİ: Arka plan set edilirken prs parametresi eklendi
                            set_slide_background(slide, temp_bg_path, prs)
                            
                            config = position_config.get(presentation_type, {})
                            
                            if img_url1:
                                add_image_to_slide(slide, img_url1, presentation_type, 0, config, image_cache)
                                if caption_columns:
                                    add_caption_to_slide(slide, row, caption_columns, presentation_type, 0, config, font_title_bold, font_title_word_wrap)
                                elif use_column_titles:
                                    add_title_to_slide(slide, col1, presentation_type, 0, config, font_title_bold, font_title_word_wrap)
                                if label_columns and 1 <= label_count: add_labels_for_image(slide, row, label_columns, 'label1', add_label_border, config)
                            if img_url2:
                                add_image_to_slide(slide, img_url2, presentation_type, 1, config, image_cache)
                                if caption_columns:
                                    add_caption_to_slide(slide, row, caption_columns, presentation_type, 1, config, font_title_bold, font_title_word_wrap)
                                elif use_column_titles:
                                    add_title_to_slide(slide, col2, presentation_type, 1, config, font_title_bold, font_title_word_wrap)
                                if label_columns and 2 <= label_count: add_labels_for_image(slide, row, label_columns, 'label2', add_label_border, config)
                else:
                    if presentation_type == 'triple' and skip_empty_rows:
                        all_image_urls = []; all_image_columns = []; all_image_row_data = []
                        for index, row in df.iterrows():
                            image_urls = []; image_cols = []
                            for col in image_columns:
                                if pd.notna(row[col]) and row[col] != '':
                                    image_urls.append(row[col]); image_cols.append(col)
                            if image_urls:
                                all_image_urls.extend(image_urls); all_image_columns.extend(image_cols); all_image_row_data.extend([row] * len(image_urls))
                        if not all_image_urls: continue
                        images_per_slide = 3
                        for i in range(0, len(all_image_urls), images_per_slide):
                            slide = prs.slides.add_slide(slide_layout)
                            # GÜNCELLENDİ: Arka plan set edilirken prs parametresi eklendi
                            set_slide_background(slide, temp_bg_path, prs)
                            
                            current_images = all_image_urls[i:i+images_per_slide]
                            current_columns = all_image_columns[i:i+images_per_slide]
                            current_row_data_list = all_image_row_data[i:i+images_per_slide]
                            config = position_config.get(presentation_type, {})
                            for j, img_url in enumerate(current_images):
                                add_image_to_slide(slide, img_url, presentation_type, j, config, image_cache)
                                if caption_columns and j < len(current_row_data_list):
                                    add_caption_to_slide(slide, current_row_data_list[j], caption_columns, presentation_type, j, config, font_title_bold, font_title_word_wrap)
                                elif use_column_titles and j < len(current_columns):
                                    column_title = current_columns[j]
                                    add_title_to_slide(slide, column_title, presentation_type, j, config, font_title_bold, font_title_word_wrap)
                                if label_columns and j < len(current_row_data_list):
                                    if (j + 1) <= label_count:
                                        current_row_data = current_row_data_list[j]
                                        label_config_key = f'label{j+1}'
                                        add_labels_for_image(slide, current_row_data, label_columns, label_config_key, add_label_border, config)
                    else:
                        for index, row in df.iterrows():
                            image_urls = []
                            for col in image_columns:
                                if pd.notna(row[col]): image_urls.append(row[col])
                            if skip_single_images and presentation_type != 'single':
                                required_images = 2 if presentation_type == 'double' else 3
                                if len(image_urls) < required_images: continue
                            if not image_urls: continue
                            if presentation_type == 'single': images_per_slide = 1
                            elif presentation_type == 'double': images_per_slide = 2
                            else: images_per_slide = 3
                            
                            for i in range(0, len(image_urls), images_per_slide):
                                slide = prs.slides.add_slide(slide_layout)
                                # GÜNCELLENDİ: Arka plan set edilirken prs parametresi eklendi
                                set_slide_background(slide, temp_bg_path, prs)
                                
                                current_images = image_urls[i:i+images_per_slide]
                                config = position_config.get(presentation_type, {})
                                for j, img_url in enumerate(current_images):
                                    add_image_to_slide(slide, img_url, presentation_type, j, config, image_cache)
                                    if caption_columns:
                                        add_caption_to_slide(slide, row, caption_columns, presentation_type, j, config, font_title_bold, font_title_word_wrap)
                                    elif use_column_titles and j < len(image_columns):
                                        column_title = image_columns[j]
                                        add_title_to_slide(slide, column_title, presentation_type, j, config, font_title_bold, font_title_word_wrap)
                                    if label_columns:
                                        if presentation_type == 'single':
                                            for k in range(label_count):
                                                label_config_key = f'label{k+1}'
                                                add_labels_for_image(slide, row, label_columns, label_config_key, add_label_border, config)
                                        elif (j + 1) <= label_count:
                                            label_config_key = f'label{j+1}'
                                            add_labels_for_image(slide, row, label_columns, label_config_key, add_label_border, config)
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        if temp_bg_path and os.path.exists(temp_bg_path): os.unlink(temp_bg_path)
        return send_file(output, as_attachment=True, download_name='presentation.pptx', mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    
    except Exception as e:
        if temp_bg_path and os.path.exists(temp_bg_path): os.unlink(temp_bg_path)
        return jsonify({'error': f'Could not generate presentation: {str(e)}'}), 500

if __name__ == '__main__':
    config.UPLOAD_ROOT.mkdir(parents=True, exist_ok=True)
    # threaded: concurrent requests. use_reloader=False avoids SSE drop on reload.
    app.run(debug=True, threaded=True, use_reloader=False)